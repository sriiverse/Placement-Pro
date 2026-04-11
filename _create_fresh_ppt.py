from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def apply_bullet_formatting(tf, bullets):
    tf.clear()
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.space_after = Pt(14)

def run():
    prs = Presentation()
    
    # Optional: adjust slide width for widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # --- 0. Title Slide ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "PlacementPro+"
    subtitle.text = "A Neuro-Symbolic AI Engine for Placement Prediction & Personalized Roadmaps\nReview Presentation"

    # --- 1. Problem Statement ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Problem Statement"
    bullets = [
        "Generic Preparation: Existing platforms offer one-size-fits-all roadmaps, ignoring base skills.",
        "Lack of Logical Progression in AI: Standard LLMs hallucinate learning paths, failing to enforce prerequisites (e.g., trying to teach ML before Python basics).",
        "Siloed Systems: A fundamental disconnect exists between ML-based prediction models and actionable GenAI curriculums."
    ]
    apply_bullet_formatting(slide.placeholders[1].text_frame, bullets)

    # --- 2. Motivations ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Motivations"
    bullets = [
        "Personalized Learning: Empower students with tailored 12-week roadmaps that react dynamically to their current skills.",
        "Precision over Hallucination: A Neuro-Symbolic approach guarantees logically sound, executable study plans without AI hallucination.",
        "Holistic Solution: Bridge statistical prediction (likelihood of placement) with intelligent recommendation (how to improve)."
    ]
    apply_bullet_formatting(slide.placeholders[1].text_frame, bullets)

    # --- 3. Summary of existing approaches and research limitations ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Existing Approaches & Limitations"
    bullets = [
        "Generative AI (LLMs): Highly flexible in text generation, but extremely prone to logical hallucinations and ignoring educational prerequisites.",
        "Traditional Recommender Systems (ML): Good at suggesting discrete courses but suffer from cold-start problems and cannot generate continuous timelines.",
        "Graph-Based Tutors (Symbolic Logic): Mathematically sound and enforces rules perfectly, but rigid and lacks the semantic understanding needed for dynamic user input."
    ]
    apply_bullet_formatting(slide.placeholders[1].text_frame, bullets)

    # --- 4. Research Gap & Improvements ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Research Gap & Improvements"
    bullets = [
        "The Gap: Current EdTech AI relies either entirely on rigid structural graphs or entirely on probabilistic language models.",
        "Our Improvement: Bridging the gap using a customized Neuro-Symbolic architecture.",
        "Neural Component (Semantic Understanding): Using sentence embeddings (SentenceTransformers) to semantically match a student's unstructured skills to the backend database.",
        "Symbolic Component (Logical Constraints): Navigating a weighted Knowledge Graph (NetworkX) to enforce strict prerequisite traversal for the generated syllabus."
    ]
    apply_bullet_formatting(slide.placeholders[1].text_frame, bullets)

    # --- 5. Proposed Solution & Architecture ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Proposed Solution & Architecture"
    bullets = [
        "Workflow Automation: User input -> ML Prediction -> NLP Vectorization -> Graph Traversal -> Syllabus Output.",
        "1. User Profiling: Frontend captures technical skills and target roles.",
        "2. Predictive Engine: Supervised ML models assess current placement likelihood score.",
        "3. Semantic Encoding: Converts strings to mathematical vectors to compute similarity.",
        "4. Knowledge Graph Matching: Deep topological traversal to fetch related topics.",
        "5. Output Generation: A formatted 12-week roadmap is delivered to the dashboard."
    ]
    apply_bullet_formatting(slide.placeholders[1].text_frame, bullets)

    # --- 6. Dataset & Tools/Technologies ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Dataset & Tools/Technologies"
    bullets = [
        "Datasets: Simulated historical student profiles (for ML scoring) and a custom Directed Acyclic Graph (DAG) for the CS Curriculum.",
        "Frontend Application: Built with React, TypeScript, and Vite for performance.",
        "Backend Architecture: FastAPI (Python) for rapid ML inference and routing.",
        "AI/ML Stack:",
        "  - SentenceTransformers (NLP Vectors)",
        "  - NetworkX (Graph Theory)",
        "  - Scikit-Learn (Predictive ML)"
    ]
    apply_bullet_formatting(slide.placeholders[1].text_frame, bullets)

    # --- 7. Progress Implementation Plan & Methodology ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Progress Implementation Plan & Methodology"
    bullets = [
        "System Skeleton: Initialized the full-stack architecture with functional React-FastAPI communication.",
        "Predictive Engine: Successfully deployed the ML-driven placement prediction component calculating readiness percentiles.",
        "Neuro-Symbolic Engine: Completely overhauled the backend API to parse skills through our custom semantic knowledge graph.",
        "Frontend Integration: Designed a dynamic, interactive dashboard to display these roadmaps continuously."
    ]
    apply_bullet_formatting(slide.placeholders[1].text_frame, bullets)

    # --- 8. Future Work Plan ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Future Work Plan"
    bullets = [
        "Advanced State Tracking: Implement granular database tracking for students' weekly progress checks.",
        "Dynamic Re-Routing: Develop the logic to auto-regenerate the remaining weeks of the roadmap if a student fails a milestone quiz.",
        "Interactive Graph Visualizations: Render the prerequisite graph visually on the frontend.",
        "Production Deployment: Containerize the FastAPI backend and deploy via standard cloud services."
    ]
    apply_bullet_formatting(slide.placeholders[1].text_frame, bullets)

    output_path = r"d:\FRP\PlacementPro_Review_Clean.pptx"
    prs.save(output_path)
    print(f"File saved to {output_path}")

if __name__ == '__main__':
    run()
