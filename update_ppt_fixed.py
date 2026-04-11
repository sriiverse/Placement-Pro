import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

def add_text_slide(prs, title, text_lines, insert_index, title_size=32, point_size=18):
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    new_idx = len(prs.slides) - 1
    
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        if shape == slide.shapes.title:
            shape.text = title
            for p in shape.text_frame.paragraphs:
                p.font.size = Pt(title_size)
        elif shape.placeholder_format.idx == 1:
            shape.text = "\n".join(text_lines)
            tf = shape.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            for p in tf.paragraphs:
                p.font.size = Pt(point_size)
            
    move_slide(prs, new_idx, insert_index)

def add_native_arch_slide(prs, title, insert_index):
    # Use blank or title-only
    try:
        slide_layout = prs.slide_layouts[5]
    except:
        slide_layout = prs.slide_layouts[1]
        
    slide = prs.slides.add_slide(slide_layout)
    new_idx = len(prs.slides) - 1
    
    if slide.shapes.title:
        slide.shapes.title.text = title
        
    shapes = slide.shapes
    
    def add_box(left, top, width, height, text, bg_rgb=(10, 30, 60)):
        box = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.text = text
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(bg_rgb[0], bg_rgb[1], bg_rgb[2])
        for p in box.text_frame.paragraphs:
            p.font.size = Pt(14)
        return box

    # Calculate positioning for absolute fit within slide
    add_box(0.5, 2.5, 2.5, 1.2, "Frontend (React / Vite)\nNetlify Global Edge CDN", (3, 7, 18))
    add_box(4.0, 2.5, 2.5, 1.2, "Backend API (Flask)\nHugging Face Containers", (24, 76, 120))
    add_box(7.5, 1.5, 2.5, 1.2, "Neuro-Symbolic Engine\nVector Embeddings + Graph", (181, 55, 242))
    add_box(7.5, 3.5, 2.5, 1.2, "Database Config\nSQLite + SQLAlchemy ORM", (44, 241, 138))

    # Add Connector Arrows
    arrow1 = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.0), Inches(2.85), Inches(1.0), Inches(0.5))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = RGBColor(100, 100, 100)
    
    arrow2 = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(2.1), Inches(1.0), Inches(0.5))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = RGBColor(100, 100, 100)

    arrow3 = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(3.5), Inches(1.0), Inches(0.5))
    arrow3.fill.solid()
    arrow3.fill.fore_color.rgb = RGBColor(100, 100, 100)

    move_slide(prs, new_idx, insert_index)

def main():
    filepath = r"d:\FRP\Frp presentation.pptx"
    prs = Presentation(filepath)
    
    intro_lines = [
        "Project Overview:",
        " • Placement.OS is a Next-Generation AI-powered Placement Preparation System.",
        " • Utilizes a Neuro-Symbolic AI architecture to map student skills dynamically.",
        " • Predicts placement probabilities and uncovers structural knowledge gaps.",
        "",
        "Why This Research Idea?",
        " • Technical education suffers from heavy academic-industry misalignment.",
        " • Students often graduate lacking the agile skills strictly required by top firms.",
        " • Our Solution: Mathematically bridge the gap by generating personalized 12-week",
        "   learning roadmaps, replacing generic and outdated educational syllabi."
    ]
    add_text_slide(prs, "Introduction & Research Motivation", intro_lines, insert_index=2, point_size=16)
    
    # Original slide 8 was at index 7. Due to intro insert, it became index 8. Insert next at index 9.
    add_native_arch_slide(prs, "Application Architecture Flow", 9)
    
    tech_lines = [
        "Frontend Architecture (Client-Side Rendering):",
        " • Core: React 19 + TypeScript (via Vite esbuild)",
        " • Aesthetics: Tailwind CSS + Framer Motion (Declarative physics)",
        " • Visualizations: Recharts (Dynamic SVG charts)",
        "",
        "Backend Architecture (Neuro-Symbolic Engine):",
        " • Framework: Python Flask (RESTful API ecosystem)",
        " • Semantic AI: HF sentence-transformers (Vector extraction)",
        " • Logic AI: NetworkX (Directed Acyclic Graph routing)",
        "",
        "Persistence & Cloud Infrastructure:",
        " • Database: SQLAlchemy Object-Relational Mapper (SQLite)",
        " • Deployment: Netlify CDN (Frontend) & HF Spaces Docker (Backend)"
    ]
    add_text_slide(prs, "Technology Stack Details", tech_lines, insert_index=10, point_size=16)
    
    # To be absolutely sure we don't mess up their base file, export to a new name
    outpath = r"d:\FRP\Frp presentation final fixed.pptx"
    prs.save(outpath)
    print("Done")

if __name__ == "__main__":
    main()
