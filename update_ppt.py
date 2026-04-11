import os
from pptx import Presentation
from pptx.util import Inches

def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

def add_text_slide(prs, title, text, insert_index):
    # Layout 1 is usually Title and Content
    slide_layout = prs.slide_layouts[1] 
    slide = prs.slides.add_slide(slide_layout)
    new_idx = len(prs.slides) - 1
    
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        if shape == slide.shapes.title:
            shape.text = title
        elif shape.placeholder_format.idx == 1:
            shape.text = text
            
    move_slide(prs, new_idx, insert_index)

def add_image_slide(prs, title, img_path, insert_index):
    # Layout 5 is usually Title Only
    try:
        slide_layout = prs.slide_layouts[5]
    except:
        slide_layout = prs.slide_layouts[1]
        
    slide = prs.slides.add_slide(slide_layout)
    new_idx = len(prs.slides) - 1
    
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Calculate rough center
    width = prs.slide_width
    height = prs.slide_height
    
    # Define max image width as 80% slide width
    max_img_width = width * 0.8
    slide.shapes.add_picture(img_path, width * 0.1, Inches(1.5), width=int(max_img_width))
    
    move_slide(prs, new_idx, insert_index)

def main():
    filepath = r"d:\FRP\Frp presentation.pptx"
    prs = Presentation(filepath)
    
    # 1. Add Introduction after Slide 2 (which is index 1). So we insert at index 2.
    intro_title = "Introduction & Research Motivation"
    intro_text = (
        "Project Overview:\n"
        "Placement.OS is an advanced AI-powered Placement Preparation System. "
        "It utilizes a Neuro-Symbolic Artificial Intelligence architecture to algorithmically "
        "evaluate student skills, predict deterministic placement probabilities, and identify specific knowledge gaps.\n\n"
        "Why This Research Idea?\n"
        "Modern technical education frequently struggles with academic-industry misalignment, where students graduate "
        "lacking the dynamic skills strictly required by top-tier tech companies. We initiated this research to mathematically "
        "bridge that gap. Rather than relying on generic educational syllabi, our system dynamically generates fully personalized, "
        "data-driven 12-week learning roadmaps."
    )
    add_text_slide(prs, intro_title, intro_text, 2)
    
    # After inserting at index 2, the original Slide 8 (index 7) has been pushed to index 8.
    # The user wants the Architecture slide AFTER the original Slide 8.
    # So we insert it at index 9.
    img_path = r"d:\FRP\docs\architecture.png"
    if os.path.exists(img_path):
        add_image_slide(prs, "Application Architecture", img_path, 9)
    else:
        print(f"Warning: Could not find image at {img_path}")
        
    # Now we add the Tech Stack slide directly after the Architecture slide (which is now at index 9).
    # So we insert at index 10.
    tech_title = "Technology Stack"
    tech_text = (
        "Frontend (Client-Side Rendering):\n"
        "• Core: React 19 + TypeScript\n"
        "• Compilation: Vite (esbuild)\n"
        "• UI & Animation: Tailwind CSS + Framer Motion\n"
        "• Visualizations: Recharts\n\n"
        "Backend (Neuro-Symbolic Engine):\n"
        "• Framework: Python Flask (RESTful API Ecosystem)\n"
        "• Vector Search: Hugging Face sentence-transformers (all-MiniLM)\n"
        "• Knowledge Graph Algorithmic Routing: NetworkX\n"
        "• Data Handling: Scikit-learn, Numpy\n\n"
        "Database Architecture:\n"
        "• SQLAlchemy mapped to SQLite\n\n"
        "Cloud Infrastructure:\n"
        "• Frontend Hosting: Netlify Global Edge CDN\n"
        "• Backend Engine: Hugging Face Spaces (Dockerized Container)"
    )
    add_text_slide(prs, tech_title, tech_text, 10)
    
    outpath = r"d:\FRP\Frp presentation final.pptx"
    prs.save(outpath)
    print(f"Success! Saved to {outpath}")

if __name__ == "__main__":
    main()
