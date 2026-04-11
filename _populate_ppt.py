from pptx import Presentation
from pptx.util import Inches, Pt
import os

template_path = r"d:\FRP\[Template] FRP Review-1 Presentation PPT .pptx"
output_path = r"d:\FRP\Final_PlacementPro_Review_Presentation.pptx"

try:
    prs = Presentation(template_path)
    
    content_map = {
        0: {
            "replaces": {
                "[Title of the Project]": "PlacementPro+: A Neuro-Symbolic AI Engine for Placement Prediction",
                "Name of the Student(s) with Regd. No.:": "Name of the Student(s): [Add Name]",
                "Group No.:": "Group No.: [Add Group]"
            }
        },
        2: { # Problem Statement
            "bullets": [
                "Generic Preparation: Existing platforms offer one-size-fits-all roadmaps, ignoring base skills.",
                "Lack of Logical Progression in AI: Standard LLMs hallucinate learning paths, failing to enforce prerequisites (e.g. ML before Python).",
                "Siloed Systems: A fundamental disconnect exists between ML-based prediction and actionable GenAI curriculums."
            ]
        },
        3: { # Motivations
            "bullets": [
                "Personalized Learning: Empower students with tailored 12-week roadmaps that react to their skills.",
                "Precision over Hallucination: A Neuro-Symbolic approach guarantees logically sound, executable study plans.",
                "Holistic Solution: Bridge statistical prediction (likelihood of placement) with intelligent recommendation (how to improve)."
            ]
        },
        5: { # Research Gap & Improvements
             "bullets": [
                 "The Gap: Current EdTech AI relies either entirely on rigid structural graphs or entirely on probabilistic language models.",
                 "Our Improvement: We use a Neuro-Symbolic architecture.",
                 "Neuro (Semantic): Sentence embeddings (SentenceTransformers) semantically match a student's unstructured skills to backend topics.",
                 "Symbolic (Logical): A weighted Knowledge Graph (NetworkX) transverses paths, ensuring the syllabus adheres to prerequisite dependencies."
             ]
        },
        6: { # Proposed Solution Architecture
             "bullets": [
                 "[INSERT ARCHITECTURE WORKFLOW DIAGRAM HERE]",
                 "1. User Profiling: Frontend captures skills/roles.",
                 "2. Predictive Engine: ML models assess placement likelihood.",
                 "3. Semantic Encoding: Encodes user skills into vectors.",
                 "4. Graph Traversal: Engine navigates a NetworkX directed graph.",
                 "5. Output: Structured 12-week roadmap is sent to the dashboard."
             ]
        },
        7: { # Proposed Solution diagram
             "bullets": [
                 "[INSERT ARCHITECTURE MODEL DIAGRAM HERE]",
                 "Frontend: React, TypeScript, and Vite.",
                 "Backend: FastAPI in Python.",
                 "AI/ML Engine: Placement Prediction & Neuro-Symbolic Roadmap Generator."
             ]
        },
        8: { # Dataset & Tools
             "bullets": [
                 "Datasets: Simulated historical student profiles, academic scores, and a custom Curriculum Graph.",
                 "Frontend Stack: React, TypeScript, Vite, Vanilla CSS.",
                 "Backend Framework: FastAPI (Python).",
                 "AI & ML Tools: SentenceTransformers, NetworkX, Scikit-Learn."
             ]
        },
        9: { # Progress
             "bullets": [
                 "System Skeleton: Successfully initialized the full-stack architecture.",
                 "Predictive Engine: Deployed the ML-driven placement prediction component.",
                 "Neuro-Symbolic Engine: Developed the backend using SentenceTransformers and NetworkX for the AI semantic graph.",
                 "Frontend Integration: Designed a dynamic dashboard to display the 12-week roadmaps natively."
             ]
        },
        10: { # Future Work
             "bullets": [
                 "Advanced State Tracking: Implement tracking for students' weekly progress.",
                 "Dynamic Re-Routing: Automatically regenerate roadmaps if a student falls behind.",
                 "Production Deployment: Prepare the application for cloud deployment using Docker."
             ]
        }
    }

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx in content_map:
            data = content_map[slide_idx]
            
            # Simple text replacement for Title slide
            if "replaces" in data:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for placeholder, new_val in data["replaces"].items():
                            if placeholder in shape.text:
                                shape.text = shape.text.replace(placeholder, new_val)
            
            # Adding bullets for content slides
            if "bullets" in data:
                # Find the shape that is the body placeholder, or just add a new textbox
                # Most templates have a title and a body placeholder. We can try to find an empty shape or just make a new one.
                # To avoid overlapping with template layout, let's create a new text box that takes up the main area.
                left = Inches(0.5)
                top = Inches(1.6)
                width = Inches(9.0)
                height = Inches(5.0)
                
                # Optionally clear old generic single-digit numbers or placeholders
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        if shape.text.strip().isdigit() and len(shape.text.strip()) < 3:
                            shape.text = ""
                
                txb = slide.shapes.add_textbox(left, top, width, height)
                tf = txb.text_frame
                tf.word_wrap = True
                
                for i, bullet in enumerate(data["bullets"]):
                    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    p.text = bullet
                    p.level = 0
                    p.font.size = Pt(20)

    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

except Exception as e:
    print(f"Error generating presentation: {e}")
